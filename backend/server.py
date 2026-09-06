import os
import uuid
import json
import base64
import asyncio
import websockets
from fastapi import FastAPI, File, UploadFile, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from pydantic import BaseModel
import uvicorn
from rag_pipeline import RAGPipeline
from dotenv import load_dotenv

# Ensure absolute paths for dotenv and frontend mounting
base_dir = os.path.dirname(os.path.abspath(__file__))
root_env = os.path.join(base_dir, "..", ".env")
backend_env = os.path.join(base_dir, ".env")

if os.path.exists(root_env):
    load_dotenv(root_env)
elif os.path.exists(backend_env):
    load_dotenv(backend_env)
else:
    load_dotenv()

frontend_dir = os.path.join(base_dir, "..", "frontend")

app = FastAPI(title="Voice RAG API - Gemini Live")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

rag = None

@app.on_event("startup")
async def startup_event():
    global rag
    print("Initializing RAG Pipeline on worker startup...")
    rag = RAGPipeline()

# Ensure absolute paths for frontend mounting
base_dir = os.path.dirname(os.path.abspath(__file__))
frontend_dir = os.path.join(base_dir, "..", "frontend")

import io

import numpy as np
from PIL import Image
import os

# Disable the new PIR API to fix the ConvertPirAttribute2RuntimeAttribute crash in recent Paddle versions
os.environ['FLAGS_enable_pir_api'] = '0'

# Initialize PaddleOCR globally so model weights load only once
# use_angle_cls=True helps with rotated text
try:
    from paddleocr import PaddleOCR
    ocr_model = PaddleOCR(
        use_angle_cls=False, # Skips the angle classification pass to save 50% time
        lang='en',
        enable_mkldnn=False,
        cpu_threads=4 # Limit threads so it doesn't fight the OS
        )
except ImportError:
    ocr_model = None
    print("PaddleOCR not installed. OCR will be disabled.")

import threading
import time

paddle_lock = threading.Lock()

def perform_ocr_hybrid(image_bytes: bytes, mime_type: str) -> str:
    from google import genai
    from google.genai import types
    
    try:
        client = genai.Client()
        for attempt in range(5):
            try:
                response = client.models.generate_content(
                    model='gemini-3.6-flash',
                    contents=[
                        types.Part.from_bytes(data=image_bytes, mime_type=mime_type),
                        "Extract all text from this accurately."
                    ]
                )
                if response and response.text:
                    return response.text
                return ""
            except Exception as e:
                print(f"Gemini OCR attempt {attempt+1} failed: {e}")
                if attempt < 4:
                    time.sleep(4)
    except Exception:
        pass
        
    print("Gemini OCR exhausted. Routing to Local OCR Cascade...")
    
    try:
        image = Image.open(io.BytesIO(image_bytes)).convert('RGB')
    except Exception:
        return ""
        
    # === Fast CPU Fallback: Tesseract ===
    try:
        import pytesseract
        import re
        import os
        
        # Configure Tesseract path if needed for Windows
        tess_path = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
        if os.path.exists(tess_path):
            pytesseract.pytesseract.tesseract_cmd = tess_path
            
        tesseract_text = pytesseract.image_to_string(image).strip()
        
        # Quality Check (Heuristic: are there enough valid words/alphanumeric chars?)
        if len(tesseract_text) > 5:
            # Count standard characters vs total
            standard_chars = len(re.findall(r'[a-zA-Z0-9\s.,!?:;\'"()-]', tesseract_text))
            ratio = standard_chars / len(tesseract_text)
            
            if ratio > 0.8: # Good quality read!
                print("Tesseract OCR succeeded with GOOD quality.")
                return tesseract_text
            else:
                print(f"Tesseract OCR quality BAD (Garbage ratio high: {ratio:.2f}). Escaping to PaddleOCR...")
        else:
            print("Tesseract OCR quality BAD (Too little text). Escaping to PaddleOCR...")
            
    except Exception as e:
        print(f"Tesseract skipped/failed ({e}). Routing directly to PaddleOCR...")

    # === Complex Fallback: PaddleOCR ===
    if not ocr_model:
        return ""
        
    try:
        with paddle_lock:
            img_array = np.array(image)
            
            result = ocr_model.ocr(img_array)
            if not result or not result[0]:
                return ""
                
            text_parts = []
            for line in result[0]:
                text_parts.append(line[1][0])
                
            print("PaddleOCR completed successfully.")
            return "\n".join(text_parts)
    except Exception as e:
        print(f"PaddleOCR fallback failed: {e}")
        return ""

def extract_pdf_text(file_bytes: bytes) -> str:
    import fitz
    import concurrent.futures
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text_parts = []
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
        futures = []
        for page in doc:
            # 1. Get standard text
            text = page.get_text()
            if text.strip():
                text_parts.append(text)
                
            # 2. Extract and OCR any embedded images on the page
            image_list = page.get_images(full=True)
            for img in image_list:
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    ext = base_image["ext"]
                    mime = "image/jpeg" if ext in ("jpeg", "jpg") else "image/png"
                    futures.append(executor.submit(perform_ocr_hybrid, image_bytes, mime))
                except Exception as e:
                    print(f"Skipping an embedded image due to error: {e}")
                    
            # 3. Fallback for completely unparseable scanned pages that hide images in weird streams
            if len(text.strip()) < 30 and not image_list:
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("png")
                futures.append(executor.submit(perform_ocr_hybrid, img_bytes, "image/png"))

        for future in futures:
            try:
                ocr_text = future.result()
                if ocr_text and ocr_text.strip():
                    text_parts.append(f"[Image Content]: {ocr_text.strip()}")
            except Exception:
                pass

    return "\n".join(text_parts)

def extract_docx_text(file_bytes: bytes) -> str:
    from docx import Document
    import concurrent.futures
    doc = Document(io.BytesIO(file_bytes))
    text_parts = []
    
    # 1. Text from paragraphs
    for p in doc.paragraphs:
        if p.text.strip():
            text_parts.append(p.text)
            
    # 2. Text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    text_parts.append(cell.text.strip())
                    
    # 3. OCR on embedded images
    with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
        futures = []
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    image_bytes = rel.target_part.blob
                    mime = getattr(rel.target_part, "content_type", "image/jpeg")
                    futures.append(executor.submit(perform_ocr_hybrid, image_bytes, mime))
                except Exception as e:
                    print(f"Skipping docx image due to error: {e}")
                    
        for future in futures:
            try:
                ocr_text = future.result()
                if ocr_text and ocr_text.strip():
                    text_parts.append(f"[Image Content]: {ocr_text.strip()}")
            except Exception:
                pass
                
    return "\n".join(text_parts)

def extract_pptx_text(file_bytes: bytes) -> str:
    from pptx import Presentation
    import concurrent.futures
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
        futures = []
        for slide in prs.slides:
            for shape in slide.shapes:
                # 1. Text from shapes
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
                    
                # 2. OCR on pictures
                if hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob
                        mime = getattr(shape.image, "content_type", "image/jpeg")
                        futures.append(executor.submit(perform_ocr_hybrid, image_bytes, mime))
                    except Exception as e:
                        print(f"Skipping pptx image due to error: {e}")
                        
        for future in futures:
            try:
                ocr_text = future.result()
                if ocr_text and ocr_text.strip():
                    text_parts.append(f"[Image Content]: {ocr_text.strip()}")
            except Exception:
                pass
                    
    return "\n".join(text_parts)

from fastapi.staticfiles import StaticFiles
import os

# Create uploads directory
UPLOADS_DIR = os.path.join(base_dir, "uploads")
os.makedirs(UPLOADS_DIR, exist_ok=True)
app.mount("/uploads", StaticFiles(directory=UPLOADS_DIR), name="uploads")

@app.post("/api/upload")
async def upload_doc(file: UploadFile = File(...)):
    content = await file.read()
    filename = file.filename.lower()
    
    text = ""
    try:
        if filename.endswith(".pdf"):
            text = extract_pdf_text(content)
        elif filename.endswith(".docx"):
            text = extract_docx_text(content)
        elif filename.endswith(".pptx"):
            text = extract_pptx_text(content)
        elif filename.endswith((".png", ".jpg", ".jpeg")):
            mime = "image/png" if filename.endswith(".png") else "image/jpeg"
            text = perform_ocr_gemini(content, mime)
        else:
            # Fallback to plain text
            text = content.decode("utf-8", errors="ignore")
    except Exception as e:
        print(f"Error parsing {filename}: {e}")
        return {"status": "error", "message": str(e)}

    doc_id = str(uuid.uuid4())
    rag.add_document(text, doc_id, file.filename)
    
    # Save file to uploads directory for viewing
    ext = os.path.splitext(file.filename)[1]
    safe_filename = f"{doc_id}{ext}"
    file_path = os.path.join(UPLOADS_DIR, safe_filename)
    with open(file_path, "wb") as f:
        f.write(content)
        
    return {"status": "success", "doc_id": doc_id, "title": file.filename, "url": f"/uploads/{safe_filename}"}

@app.get("/api/documents")
async def get_documents():
    if rag:
        docs = rag.get_all_documents()
        # Attach URLs based on doc_id
        for doc in docs:
            # Try to find the file in uploads with the doc_id prefix
            doc["url"] = ""
            for f in os.listdir(UPLOADS_DIR):
                if f.startswith(doc["doc_id"]):
                    doc["url"] = f"/uploads/{f}"
                    break
        return {"status": "success", "documents": docs}
    return {"status": "error", "message": "RAG pipeline not initialized"}

@app.delete("/api/documents/{doc_id}")
async def delete_doc(doc_id: str):
    if rag:
        try:
            rag.delete_document(doc_id)
            # Remove file from disk
            for f in os.listdir(UPLOADS_DIR):
                if f.startswith(doc_id):
                    os.remove(os.path.join(UPLOADS_DIR, f))
                    break
            return {"status": "success"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    return {"status": "error", "message": "RAG pipeline not initialized"}

@app.websocket("/ws/live")
async def websocket_endpoint(websocket: WebSocket):
    await websocket.accept()
    
    raw_key = os.getenv("GEMINI_API_KEY", "")
    api_key = raw_key.strip().strip('"').strip("'")
    if not api_key:
        print("Error: GEMINI_API_KEY not set or empty in .env")
        await websocket.close(code=1008)
        return

    print(f"Connecting to Gemini Live WS using key starting with '{api_key[:6]}...'")

    # Both Standard (AIzaSy...) and Auth (AQ....) keys authenticate the same way
    # for the Live API WebSocket: as a `key` query parameter. Do NOT send AQ. keys
    # as an Authorization: Bearer header -- they are API keys, not OAuth access
    # tokens, and Google will reject a Bearer-style AQ. key with 1008 policy
    # violation / "Expected OAuth 2 access token...".
    gemini_ws_url = f"wss://generativelanguage.googleapis.com/ws/google.ai.generativelanguage.v1alpha.GenerativeService.BidiGenerateContent?key={api_key}"
    headers = None
    
    setup_msg = {
        "setup": {
            "model": "models/gemini-3.1-flash-live-preview",
            "generationConfig": {
                "responseModalities": ["AUDIO"]
            },
            "systemInstruction": {
                "parts": [{"text": "You are a helpful voice assistant. You can search documents to answer user queries using the search_documents tool. ALWAYS provide concise responses."}]
            },
            "tools": [{"functionDeclarations": [{
                "name": "search_documents",
                "description": "Search uploaded documents for information.",
                "parameters": {
                    "type": "OBJECT",
                    "properties": {
                        "query": {"type": "STRING", "description": "The search query"}
                    },
                    "required": ["query"]
                }
            }]}]
        }
    }
    
    try:
        connect_kwargs = {"additional_headers": headers} if headers else {}
        async with websockets.connect(gemini_ws_url, **connect_kwargs) as gemini_ws:
            await gemini_ws.send(json.dumps(setup_msg))
            
            # Wait for setup complete
            setup_resp = await gemini_ws.recv()
            print("Gemini Setup Response:", setup_resp)
            
            async def client_to_gemini():
                try:
                    while True:
                        data = await websocket.receive()
                        if 'bytes' in data:
                            msg = {
                                "realtimeInput": {
                                    "audio": {
                                        "mimeType": "audio/pcm;rate=16000", 
                                        "data": base64.b64encode(data['bytes']).decode()
                                    }
                                }
                            }
                            await gemini_ws.send(json.dumps(msg))
                        elif 'text' in data:
                            msg = json.loads(data['text'])
                            if msg.get('type') == 'client_content':
                                content_msg = {
                                    "clientContent": {
                                        "turns": [{
                                            "role": "user",
                                            "parts": [{"text": msg.get('text', '')}]
                                        }],
                                        "turnComplete": True
                                    }
                                }
                                await gemini_ws.send(json.dumps(content_msg))
                except WebSocketDisconnect:
                    pass
                except Exception as e:
                    print("Error in client_to_gemini:", e)
                    
            async def gemini_to_client():
                try:
                    while True:
                        resp = await gemini_ws.recv()
                        msg = json.loads(resp)
                        
                        if "serverContent" in msg:
                            model_turn = msg["serverContent"].get("modelTurn", {})
                            for part in model_turn.get("parts", []):
                                if "inlineData" in part:
                                    audio_data = base64.b64decode(part["inlineData"]["data"])
                                    await websocket.send_bytes(audio_data)
                        elif "toolCall" in msg:
                            calls = msg["toolCall"]["functionCalls"]
                            responses = []
                            for call in calls:
                                if call["name"] == "search_documents":
                                    query = call.get("args", {}).get("query", "")
                                    print(f"Tool call: search_documents(query='{query}')")
                                    docs = rag.retrieve(query, top_k=5)
                                    context = "\n".join(docs) if docs else "No documents found."
                                    responses.append({
                                        "id": call["id"],
                                        "name": call["name"],
                                        "response": {"result": context}
                                    })
                            if responses:
                                tool_resp = {
                                    "toolResponse": {
                                        "functionResponses": responses
                                    }
                                }
                                await gemini_ws.send(json.dumps(tool_resp))
                except Exception as e:
                    print("Error in gemini_to_client:", e)

            await asyncio.gather(client_to_gemini(), gemini_to_client())
    except Exception as e:
        print("Gemini WS Error:", e)
    finally:
        try:
            await websocket.close()
        except:
            pass

app.mount("/static", StaticFiles(directory=frontend_dir), name="static")

@app.get("/4c277e81f6fa23aec2a198d848d3ef6c.gif")
async def get_orb_gif():
    return FileResponse(os.path.join(frontend_dir, "4c277e81f6fa23aec2a198d848d3ef6c.gif"))

@app.get("/")
async def root():
    return FileResponse(os.path.join(frontend_dir, "index.html"))

if __name__ == "__main__":
    uvicorn.run("server:app", host="0.0.0.0", port=8000, reload=True)